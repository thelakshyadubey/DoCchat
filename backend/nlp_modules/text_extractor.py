import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
import os

async def extract_text_from_file(file_path: str, file_extension: str) -> str:
    """
    Extract text from various file formats
    """
    text = ""
    
    try:
        if file_extension.lower() == "pdf":
            text = extract_from_pdf(file_path)
        elif file_extension.lower() in ["docx", "doc"]:
            text = extract_from_docx(file_path)
        elif file_extension.lower() in ["pptx", "ppt"]:
            text = extract_from_pptx(file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_extension}")
        
        return text.strip()
    except Exception as e:
        raise Exception(f"Error extracting text: {str(e)}")

def extract_from_pdf(file_path: str) -> str:
    """Extract text from PDF using PyMuPDF"""
    text = ""
    try:
        doc = fitz.open(file_path)
        for page in doc:
            text += page.get_text()
        doc.close()
    except Exception as e:
        raise Exception(f"PDF extraction error: {str(e)}")
    
    return text

def extract_from_docx(file_path: str) -> str:
    """Extract text from DOCX files"""
    text = ""
    try:
        doc = Document(file_path)
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
    except Exception as e:
        raise Exception(f"DOCX extraction error: {str(e)}")
    
    return text

def extract_from_pptx(file_path: str) -> str:
    """Extract text from PPTX files"""
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        raise Exception(f"PPTX extraction error: {str(e)}")
    
    return text
